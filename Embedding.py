import os
import re
import faiss  # Obligatoire pour read_index, write_index, IndexFlatIP
import hashlib
import logging
import torch
import numpy as np
import faiss
from pymongo import MongoClient
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as WordDocument
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import threading
from typing import List
import datetime
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openai import OpenAI
from dotenv import load_dotenv

# Charger les variables d'environnement
load_dotenv()

pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class Embedding:
    def __init__(
        self,
        directory: str = "Source Files",
        model_name: str = "text-embedding-3-large",  # Modèle OpenAI
        chunk_size: int = 150,
        mongo_uri: str = "mongodb://localhost:27017/",
        db_name: str = "embeddings_db",
        collection_name: str = "files",
        faiss_index_path: str = "faiss_index.index",
        force_reindex=False,
        use_gpu: bool = True,
        embedding_dimensions: int = 1536  # Dimension par défaut pour text-embedding-3-small
    ):
        self.directory = os.path.abspath(directory)
        self.faiss_index_path = os.path.abspath(faiss_index_path)
        self.use_gpu = use_gpu
        self.model_name = model_name
        self.embedding_dimensions = embedding_dimensions
        
        # Configuration du GPU pour FAISS seulement (OpenAI gère ses propres ressources)
        self.device = "cuda" if torch.cuda.is_available() and self.use_gpu else "cpu"
        
        # Initialisation du client OpenAI
        self.openai_client = OpenAI(
            api_key=os.getenv('OPENAI_API_KEY')
        )
        
        # Vérifier que la clé API est présente
        if not os.getenv('OPENAI_API_KEY'):
            raise ValueError("Clé API OpenAI manquante. Vérifiez votre fichier .env")
        
        self.client = MongoClient(mongo_uri)
        self.db = self.client[db_name]
        self.collection = self.db[collection_name]
        self.dimension = embedding_dimensions
        self.chunk_size = chunk_size

        # Initialisation FAISS avec normalisation L2
        self.index = self._init_faiss_index()

        if not os.path.exists(self.directory):
            os.makedirs(self.directory, exist_ok=True)
            logger.info(f"Répertoire créé : {self.directory}")



        self._process_existing_files()
        self._start_file_monitoring()
        
        logger.info(f"Configuration initialisée - Modèle: {self.model_name}, Dimensions: {self.dimension}")

    def _init_faiss_index(self):
        """Initialise FAISS sur CPU uniquement"""
        if os.path.exists(self.faiss_index_path):
            print("📥 Chargement de l'index FAISS existant")
            index = faiss.read_index(self.faiss_index_path)
        else:
            print(f"🆕 Création d'un nouvel index FAISS avec {self.dimension} dimensions")
            index = faiss.IndexFlatIP(self.dimension)
            faiss.write_index(index, self.faiss_index_path)

        # ❌ PAS DE GPU SUR WINDOWS → NE RIEN FAIRE ICI
        return index

    def _file_hash(self, file_path: str) -> str:
        """Calcule le hash SHA-256 d'un fichier"""
        hasher = hashlib.sha256()
        with open(file_path, "rb") as f:
            while chunk := f.read(4096):
                hasher.update(chunk)
        return hasher.hexdigest()

    def _is_file_processed(self, file_path: str) -> bool:
        """Vérifie si le fichier a déjà été traité avec le modèle actuel"""
        return self.collection.count_documents({
            "file_path": file_path,
            "file_hash": self._file_hash(file_path),
            "model_name": self.model_name  # Vérifier aussi le modèle
        }) > 0

    def _process_existing_files(self):
        """Traite tous les fichiers existants dans le dossier"""
        logger.info("Début du traitement des fichiers existants...")
        for filename in os.listdir(self.directory):
            file_path = os.path.join(self.directory, filename)
            if os.path.isfile(file_path):
                self._process_file(file_path)

    def _process_file(self, file_path: str):
        """Pipeline complet de traitement d'un fichier"""
        if self._is_file_processed(file_path):
            logger.debug(f"Fichier déjà traité : {file_path}")
            return

        try:
            pages = self._extract_text(file_path)
            if not pages:
                logger.warning(f"Fichier vide ou non textuel : {file_path}")
                return

            # Générer embeddings avec OpenAI
            embeddings = self._generate_embeddings(pages)
            self._save_to_database(file_path, pages, embeddings)
            
            logger.info(f"Fichier traité avec succès : {file_path} ({len(pages)} pages)")

        except Exception as e:
            logger.error(f"Erreur lors du traitement de {file_path} : {str(e)}")

    def _extract_text(self, file_path: str) -> List[str]:
        """Extraction du texte selon le format du fichier, retourne une liste de pages/slides"""
        try:
            if file_path.endswith('.pdf'):
                return self._extract_pdf(file_path)
            elif file_path.endswith('.pptx'):
                return self._extract_ppt(file_path)
            elif file_path.endswith('.docx'):
                return self._extract_word(file_path)
            elif file_path.endswith(('.png', '.jpeg', '.jpg')):
                text = self._extract_image(file_path)
                return [text] if text else []
            else:
                logger.warning(f"Format non supporté : {file_path}")
                return []
        except Exception as e:
            logger.error(f"Erreur d'extraction : {file_path} - {str(e)}")
            return []
        
    def _extract_image(self, file_path: str) -> str:
        try:
            img = Image.open(file_path)
            text = pytesseract.image_to_string(img, lang='fra+eng')
            return text
        except Exception as e:
            logger.error(f"Erreur lors de l'extraction du texte de l'image {file_path} : {e}")
        return ""

    def reindex_all(self):
        """Réindexe tous les fichiers (FAISS + MongoDB)"""
        logger.info("Réindexation complète...")

        # Supprimer anciens documents
        self.collection.delete_many({})
        
        # Supprimer index FAISS
        self.index.reset()
        if os.path.exists(self.faiss_index_path):
            os.remove(self.faiss_index_path)
        
        # Recréer un index vide
        self.index = self._init_faiss_index()

        # Traiter les fichiers à nouveau
        self._process_existing_files()

        # Sauvegarder FAISS
        faiss.write_index(self.index, self.faiss_index_path)

        logger.info("Réindexation terminée.")

    def _extract_pdf(self, file_path: str) -> List[str]:
        """Extraction texte PDF avec retour d'une liste de pages"""
        pages_text = []
        
        try:
            with open(file_path, "rb") as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        page_text = re.sub(r'\s+', ' ', page_text)
                        pages_text.append(page_text)
            
            if not pages_text or sum(len(p) for p in pages_text) < 50:
                pages_text = []
                images = convert_from_path(file_path)
                
                for image in images:
                    page_text = pytesseract.image_to_string(image)
                    page_text = re.sub(r'\s+', ' ', page_text)
                    pages_text.append(page_text)
        
        except Exception as e:
            print(f"Erreur lors de l'extraction du PDF: {e}")
        
        return pages_text
    
    def _extract_ppt(self, file_path: str) -> List[str]:
        """Extraction texte PowerPoint par slide"""
        prs = Presentation(file_path)
        slides_text = []

        for slide in prs.slides:
            slide_content = []

            for shape in slide.shapes:
                # Texte direct de la forme
                text = getattr(shape, "text", None)
                if text and text.strip():
                    slide_content.append(text.strip())

                # Tableaux dans les formes
                if hasattr(shape, "has_table") and shape.has_table:
                    table = getattr(shape, "table", None)
                    if table:
                        for row in table.rows:
                            for cell in row.cells:
                                cell_text = getattr(cell, "text", "")
                                if cell_text.strip():
                                    slide_content.append(cell_text.strip())

                # Formes groupées (sous-formes)
                if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    if hasattr(shape, "shapes"): 
                        for sub_shape in shape.shapes: # type: ignore[attr-defined]
                            sub_text = getattr(sub_shape, "text", None)
                            if sub_text and sub_text.strip():
                                slide_content.append(sub_text.strip())

            # Notes de la slide
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                notes_slide = getattr(slide, "notes_slide", None)
                if notes_slide and hasattr(notes_slide, "notes_text_frame"):
                    notes_text = getattr(notes_slide.notes_text_frame, "text", "")
                    if notes_text.strip():
                        slide_content.append(notes_text.strip())

            # Ajout du contenu final de la slide
            if slide_content:
                slides_text.append(" ".join(slide_content))

        return slides_text

    def _extract_word(self, file_path: str) -> List[str]:
        """Extraction texte Word"""
        doc = WordDocument(file_path)
        full_text = "\n".join(
            paragraph.text 
            for paragraph in doc.paragraphs 
            if paragraph.text.strip()
        )
        return [full_text]

    def _chunk_text(self, text: str) -> List[str]:
        max_chunk_size = self.chunk_size  # ex: 150 tokens ≈ 500-600 caractères
        sentences = re.split(r'(?<=[.!?]) +', text)

        chunks = []
        current_chunk = ""

        for sentence in sentences:
            if len(current_chunk) + len(sentence) <= max_chunk_size:
                current_chunk += " " + sentence
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = sentence

        if current_chunk:
            chunks.append(current_chunk.strip())

        return chunks

    def _generate_embeddings(self, chunks: List[str]) -> np.ndarray:
        """Génération des embeddings avec l'API OpenAI"""
        embeddings = []
        
        # Traitement par batch pour optimiser les appels API
        batch_size = 100  # OpenAI recommande des batches de 100 max
        
        for i in range(0, len(chunks), batch_size):
            batch = chunks[i:i + batch_size]
            
            try:
                response = self.openai_client.embeddings.create(
                    input=batch,
                    model=self.model_name,
                    dimensions=self.embedding_dimensions  # Spécifier les dimensions
                )
                
                batch_embeddings = [data.embedding for data in response.data]
                embeddings.extend(batch_embeddings)
                
                logger.info(f"Embeddings générés pour batch {i//batch_size + 1}/{(len(chunks) + batch_size - 1)//batch_size}")
                
            except Exception as e:
                logger.error(f"Erreur lors de la génération d'embeddings: {str(e)}")
                raise e
        
        # Convertir en numpy array et normaliser pour la similarité cosinus
        embeddings_array = np.array(embeddings, dtype=np.float32)
        
        # Normalisation L2 pour la similarité cosinus
        norms = np.linalg.norm(embeddings_array, axis=1, keepdims=True)
        norms = np.where(norms == 0, 1, norms)  # Éviter la division par zéro
        embeddings_normalized = embeddings_array / norms
        
        return embeddings_normalized

    def _save_to_database(self, file_path: str, chunks: List[str], embeddings: np.ndarray):
        """Sauvegarde dans MongoDB et FAISS"""
        start_idx = self.index.ntotal
        self.index.add(embeddings)

        document = {
            "file_path": file_path,
            "file_name": os.path.basename(file_path),
            "file_hash": self._file_hash(file_path),
            "model_name": self.model_name,  # Stocker le modèle utilisé
            "embedding_dimensions": self.embedding_dimensions,  # Stocker les dimensions
            "chunks": chunks,
            "embedding_ids": list(range(start_idx, start_idx + len(chunks))),
            "processing_date": datetime.datetime.now(),
            "chunk_count": len(chunks)
        }

        self.collection.insert_one(document)
        faiss.write_index(self.index, self.faiss_index_path)

    def _start_file_monitoring(self):
        """Surveillance automatique du dossier"""
        event_handler = FileSystemEventHandler()
        event_handler.on_created = lambda event: self._on_file_created(event)
        self.observer = Observer()
        self.observer.schedule(event_handler, self.directory, recursive=True)
        threading.Thread(target=self.observer.start, daemon=True).start()
        logger.info(f"Surveillance activée : {self.directory}")

    def _on_file_created(self, event):
        """Gestion des nouveaux fichiers"""
        if not event.is_directory:
            self._process_file(event.src_path)

    def close(self):
        """Nettoyage des ressources"""
        self.observer.stop()
        self.observer.join()
        self.client.close()
        faiss.write_index(self.index, self.faiss_index_path)
        logger.info("Système d'embedding arrêté")

    def __repr__(self):
        return (
            f"<EmbeddingSystem: "
            f"Modèle={self.model_name} "
            f"Dimensions={self.dimension} "
            f"Fichiers traités={self.collection.count_documents({})} "
            f"Chunks indexés={self.index.ntotal}>"
        )
